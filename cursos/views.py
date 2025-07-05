from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.forms import UserCreationForm, AuthenticationForm
from django.contrib.auth.models import User
from django.contrib.auth import login, logout, authenticate
from django.db import IntegrityError
from .forms import FormularioRegistro, ExamenForm
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from django.contrib.auth.decorators import login_required
from django.http import FileResponse
from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from django.utils import timezone
import requests
from django.core.files.base import ContentFile
from .models import Curso, ProgresoCurso, PerfilUsuario
import os
from io import BytesIO
from django.db import transaction
import tempfile
import urllib.request
from django.conf import settings
import os
from django.forms import modelform_factory, inlineformset_factory
from django.shortcuts import get_object_or_404, redirect, render
import stripe
from django.contrib.staticfiles import finders
from django.conf import settings
from .models import PerfilUsuario
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation as PPTXPresentation
import os
from django.conf import settings
from django.http import FileResponse
from django.core.exceptions import ObjectDoesNotExist
import boto3
from botocore.exceptions import ClientError
from django.conf import settings
from django.core.files.storage import default_storage
from django.contrib.auth.decorators import login_required, user_passes_test
from .models import Curso, ArchivoCurso
from .forms import CursoForm, ArchivoCursoFormSet, PreguntaFormSet, OpcionFormSet
from .models import CursoComprado
from .models import Curso, Examen, PreguntaExamen
from .forms import PreguntaExamenForm, OpcionRespuestaFormSet
from django.http import JsonResponse
from django.views.decorators.http import require_POST
from django.contrib.auth.decorators import login_required
from django.utils import timezone
from django.forms import formset_factory
from django.contrib import messages
from .models import ProgresoCurso, IntentoExamen, OpcionRespuesta, Examen
import os
from io import BytesIO
from django.core.files.base import ContentFile
from django.conf import settings
from django.core.files.base import ContentFile
from reportlab.lib.utils import ImageReader
import os
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from io import BytesIO
import os

stripe.api_key = settings.STRIPE_SECRET_KEY

# Create your views here.


def signup(request):
    if request.method == 'GET':
        return render(request, 'signup.html', {
            'form': FormularioRegistro()
        })

    form = FormularioRegistro(request.POST)
    if form.is_valid():
        username = form.cleaned_data['username'].strip()
        password = form.cleaned_data['password1']
        curp = form.cleaned_data['curp'].strip().upper()

        print("📥 Datos recibidos del formulario:")
        print("Nombre:", form.cleaned_data['nombre'])
        print("Apellido paterno:", form.cleaned_data['apellido_paterno'])
        print("Apellido materno:", form.cleaned_data['apellido_materno'])
        print("CURP:", curp)
        print("Puesto:", form.cleaned_data['puesto'])
        print("ocupacion_especifica:",
              form.cleaned_data['ocupacion_especifica'])
        print("nombre_razon_social:", form.cleaned_data['nombre_razon_social'])
        print("refc_empresa:", form.cleaned_data['refc_empresa'])

        # Validar duplicados antes de crear
        if User.objects.filter(username=username).exists():
            form.add_error('username', 'El nombre de usuario ya existe.')
            return render(request, 'signup.html', {'form': form})

        if PerfilUsuario.objects.filter(curp__iexact=curp).exists():
            form.add_error('curp', 'El CURP ya está registrado.')
            return render(request, 'signup.html', {'form': form})

        try:
            with transaction.atomic():  # Asegura que se crea usuario y perfil juntos
                # Crear usuario
                user = User.objects.create_user(
                    username=username, password=password)

                # Crear perfil manualmente
                perfil = PerfilUsuario.objects.create(
                    user=user,
                    nombre=form.cleaned_data['nombre'],
                    apellido_paterno=form.cleaned_data['apellido_paterno'],
                    apellido_materno=form.cleaned_data['apellido_materno'],
                    curp=curp,
                    puesto=form.cleaned_data['puesto'],
                    ocupacion_especifica=form.cleaned_data['ocupacion_especifica'],
                    nombre_razon_social=form.cleaned_data['nombre_razon_social'],
                    refc_empresa=form.cleaned_data['refc_empresa']
                )

            # Si todo salió bien
            login(request, user)
            messages.success(request, f"✅ Perfil creado para: {user.username}")
            return redirect('home')

        except IntegrityError as e:
            print("❌ INTEGRITY ERROR:", str(e))
            form.add_error(
                None, "Error al crear el usuario o perfil, posiblemente datos duplicados.")
            return render(request, 'signup.html', {'form': form})

    # Si el formulario no es válido
    return render(request, 'signup.html', {
        'form': form,
        'error': 'Formulario inválido.'
    })


def signin(request):

    if request.method == 'GET':
        return render(request, 'signin.html', {
            'form': AuthenticationForm
        })
    else:
        user = authenticate(
            request, username=request.POST['username'], password=request.POST['password'])
        if user is None:
            return render(request, 'signin.html', {
                'form': AuthenticationForm,
                'error': 'Username or password is incorect'
            })
        else:
            login(request, user)
            return redirect('home')


@login_required
def signout(request):
    logout(request)
    return redirect('home')


def convert_pptx_to_pdf(pptx_path, pdf_path):
    ppt = PPTXPresentation(pptx_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    for slide in ppt.slides:
        y_position = height - 50  # Margen superior

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines = shape.text.strip().split('\n')
                for line in lines:
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    c.drawString(50, y_position, line)
                    y_position -= 20

        c.showPage()  # Nueva página para siguiente diapositiva

    c.save()


@login_required
def home(request):

    # Debug conexión a S3 (opcional)
    try:
        s3 = boto3.client(
            's3',
            aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
            aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
            region_name=settings.AWS_S3_REGION_NAME,
        )
        response = s3.list_objects_v2(Bucket=settings.AWS_STORAGE_BUCKET_NAME)
        print("✅ Conexión S3 exitosa. Objetos encontrados:",
              response.get('KeyCount', 0))
    except ClientError as e:
        print("❌ Error de conexión con S3:", e)

    if hasattr(default_storage, 'bucket'):
        print("🔍 Archivos en S3 con prefijo 'media/':")
        for obj in default_storage.bucket.objects.filter(Prefix='media/'):
            print(" -", obj.key)

    # Listar cursos ordenados más recientes
    cursos = Curso.objects.all().order_by('-creado_en')

    cursos_con_acceso = set(
        CursoComprado.objects.filter(
            usuario=request.user, pagado=True).values_list('curso_id', flat=True)
    )

    # Renderizar plantilla de cursos
    return render(request, 'home.html', {
        'cursos': cursos,
        'cursos_con_acceso': cursos_con_acceso,
    })


def root_redirect(request):
    if request.user.is_authenticated:
        return redirect('home')
    # Renderizamos el signin con el formulario para no usar redirección 302
    return render(request, 'signin.html', {
        'form': AuthenticationForm()
    })


@login_required
def pago_cancelado_curso(request, curso_id):
    return render(request, 'pago_cancelado.html', {'curso_id': curso_id})


def es_superusuario(user):
    return user.is_superuser


@login_required
def subir_curso(request):
    if request.method == 'POST':
        curso_form = CursoForm(request.POST, request.FILES)  # <--- aquí
        formset = ArchivoCursoFormSet(
            request.POST, request.FILES, queryset=ArchivoCurso.objects.none())
        if curso_form.is_valid() and formset.is_valid():
            curso = curso_form.save(commit=False)
            curso.creado_por = request.user
            curso.save()

            for form in formset:
                if form.cleaned_data:
                    archivo = form.save(commit=False)
                    archivo.curso = curso
                    archivo.save()

            return redirect('detalle_curso', pk=curso.id)
    else:
        curso_form = CursoForm()
        formset = ArchivoCursoFormSet(queryset=ArchivoCurso.objects.none())

    return render(request, 'subir_curso.html', {
        'curso_form': curso_form,
        'formset': formset,
    })


@user_passes_test(es_superusuario)
def eliminar_archivo(request, archivo_id):
    archivo = get_object_or_404(ArchivoCurso, id=archivo_id)
    curso_id = archivo.curso.id
    archivo.delete()
    return redirect('detalle_curso', pk=curso_id)


@user_passes_test(es_superusuario)
def eliminar_curso(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)

    # Ejecuta el delete personalizado de cada archivo
    for archivo in curso.archivos.all():
        archivo.delete()

    # Luego elimina el curso
    curso.delete()

    return redirect('home')


@login_required
def crear_checkout_por_curso(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)

    # Guardar en metadata para luego asociar
    session = stripe.checkout.Session.create(
        payment_method_types=['card'],
        line_items=[{
            'price_data': {
                'currency': 'mxn',
                'product_data': {
                    'name': curso.titulo,
                },
                'unit_amount': int(curso.precio * 100),  # en centavos
            },
            'quantity': 1,
        }],
        mode='payment',
        success_url=request.build_absolute_uri(
            f'/pago_exitoso_curso/{curso.id}/'),
        cancel_url=request.build_absolute_uri(
            f'/pago_cancelado_curso/{curso.id}/'),
        metadata={
            'user_id': request.user.id,
            'curso_id': curso.id,
        }
    )
    return redirect(session.url, code=303)


@login_required
def pago_exitoso_curso(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)
    obj, creado = CursoComprado.objects.get_or_create(
        usuario=request.user,
        curso=curso
    )
    obj.pagado = True
    obj.save()
    return redirect('detalle_curso', pk=curso.id)


@login_required
def detalle_curso(request, pk):
    curso = get_object_or_404(Curso, pk=pk)

    # Validación compra para usuarios normales
    if not request.user.is_superuser:
        try:
            compra = CursoComprado.objects.get(
                usuario=request.user, curso=curso)
            if not compra.pagado:
                return redirect('crear_checkout_por_curso', curso_id=pk)
        except CursoComprado.DoesNotExist:
            return redirect('crear_checkout_por_curso', curso_id=pk)

    progreso, _ = ProgresoCurso.objects.get_or_create(
        usuario=request.user, curso=curso)

    # Si el video está marcado como completo pero examen NO aprobado, redirigir al examen
    if progreso.video_visto_completo and not progreso.examen_aprobado:
        return redirect('examen_curso', curso_id=pk)

    archivos = ArchivoCurso.objects.filter(
        curso=curso) if progreso.examen_aprobado or request.user.is_superuser else []

    formset = ArchivoCursoFormSet(
        queryset=ArchivoCurso.objects.none()) if request.user.is_superuser else None

    return render(request, 'detalle_curso.html', {
        'curso': curso,
        'archivos': archivos,
        'formset': formset,
        'progreso': progreso,
    })


@login_required
@require_POST
def marcar_video_completo(request, curso_id):
    curso = get_object_or_404(Curso, pk=curso_id)
    progreso, _ = ProgresoCurso.objects.get_or_create(
        usuario=request.user, curso=curso)

    progreso.video_visto_completo = True

    # Reinicia intentos si aún no aprobó
    if not progreso.examen_aprobado:
        IntentoExamen.objects.filter(
            usuario=request.user, examen=curso.examen).delete()

    progreso.save()

    return JsonResponse({'status': 'ok'})


@login_required
def examen_curso(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)
    examen = getattr(curso, 'examen', None)
    if not examen:
        messages.error(request, "Este curso no tiene examen asignado.")
        return redirect('detalle_curso', pk=curso_id)

    progreso, _ = ProgresoCurso.objects.get_or_create(
        usuario=request.user, curso=curso)

    # Solo puede hacer examen si video visto completo
    if not progreso.video_visto_completo:
        messages.warning(
            request, "Debes ver el video completo antes de hacer el examen.")
        return redirect('detalle_curso', pk=curso_id)

    # Checar intentos finalizados
    intentos_finalizados = IntentoExamen.objects.filter(
        usuario=request.user, examen=examen, estado='finalizado').count()
    if intentos_finalizados >= 3 and not progreso.examen_aprobado:
        messages.error(
            request, "Has agotado los 3 intentos para este examen. Debes volver a ver el video.")
        progreso.video_visto_completo = False
        progreso.save()
        return redirect('detalle_curso', pk=curso_id)

    # Ver si ya hay un intento en curso
    intento_en_curso = IntentoExamen.objects.filter(
        usuario=request.user, examen=examen, estado='en_curso').first()
    if intento_en_curso:
        # Aquí podrías implementar lógica para controlar tiempo restante, cancelar si salio de la página, etc.
        pass
    else:
        # Crear nuevo intento
        intento_en_curso = IntentoExamen.objects.create(
            usuario=request.user, examen=examen)

    preguntas = examen.preguntas.all()

    if request.method == 'POST':
        total = preguntas.count()
        correctas = 0

        for pregunta in preguntas:
            respuesta_id = request.POST.get(f'pregunta_{pregunta.id}')
            if respuesta_id:
                opcion = OpcionRespuesta.objects.filter(
                    id=respuesta_id, pregunta=pregunta).first()
                if opcion and opcion.es_correcta:
                    correctas += 1

        porcentaje = (correctas / total) * 100 if total > 0 else 0
        aprobado = porcentaje >= 90

        intento_en_curso.porcentaje = porcentaje
        intento_en_curso.aprobado = aprobado
        intento_en_curso.estado = 'finalizado'
        intento_en_curso.fecha_fin = timezone.now()
        intento_en_curso.save()

        if aprobado:
            progreso.examen_aprobado = True
            progreso.fecha_examen_aprobado = timezone.now()  # <-- nuevo
            progreso.porcentaje_examen = porcentaje           # <-- nuevo
            progreso.save()
            messages.success(request, "¡Felicidades! Has aprobado el examen.")
            return redirect('certificado_reconocimiento', curso_id=curso_id)
        else:
            messages.error(
                request, f"No alcanzaste el 90%. Tu puntaje fue {porcentaje:.2f}%. Inténtalo nuevamente.")
            return redirect('detalle_curso', pk=curso_id)

    return render(request, 'examen_curso.html', {
        'curso': curso,
        'examen': examen,
        'preguntas': preguntas,
        'intento': intento_en_curso,
        'tiempo_limite': examen.tiempo_minutos * 60,  # en segundos
    })


@login_required
def certificado_reconocimiento(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)
    progreso = get_object_or_404(ProgresoCurso, usuario=request.user, curso=curso)

    if not progreso.examen_aprobado:
        messages.error(request, "No has aprobado el examen aún.")
        return redirect('detalle_curso', pk=curso_id)

    if progreso.certificado:
        progreso.certificado.delete()

    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=landscape(letter))
    width, height = landscape(letter)

    # --- Fondo blanco ---
    c.setFillColor(colors.white)
    c.rect(0, 0, width, height, fill=1)

    # --- Imagen de fondo ---
    try:
        if settings.DEBUG:
            local_path = os.path.join(settings.BASE_DIR, 'cursos', 'static', 'logo_certificado.jpg')
            if os.path.exists(local_path):
                c.drawImage(local_path, 0, 0, width=width, height=height, preserveAspectRatio=False, mask='auto')
        else:
            s3_url = 'https://cursos-seguridad-media.s3.us-east-2.amazonaws.com/static/cursos_seguridad/img/dc3_template.jpg'
            response = requests.get(s3_url)
            if response.status_code == 200:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp_img:
                    tmp_img.write(response.content)
                    tmp_img.flush()
                    c.drawImage(tmp_img.name, 0, 0, width=width, height=height, preserveAspectRatio=False, mask='auto')
    except Exception as e:
        print(f"[Certificado] No se pudo cargar la imagen de fondo: {e}")

    def draw_centered_text(canvas, text, y, font_name, font_size, color=colors.black):
        canvas.setFont(font_name, font_size)
        canvas.setFillColor(color)
        text_width = canvas.stringWidth(text, font_name, font_size)
        x = (width - text_width) / 2
        canvas.drawString(x, y, text)

    perfil = get_object_or_404(PerfilUsuario, user=request.user)
    nombre_completo = f"{perfil.nombre} {perfil.apellido_paterno} {perfil.apellido_materno}".strip()

    nombre_curso = curso.titulo
    examen = getattr(curso, 'examen', None)
    duracion = examen.tiempo_minutos // 60 if examen and examen.tiempo_minutos else 10
    fecha_curso = progreso.fecha_inicio.strftime("%d de %B del %Y") if progreso.fecha_inicio else timezone.now().strftime("%d de %B del %Y")

    meses = {
        'January': 'Enero', 'February': 'Febrero', 'March': 'Marzo',
        'April': 'Abril', 'May': 'Mayo', 'June': 'Junio',
        'July': 'Julio', 'August': 'Agosto', 'September': 'Septiembre',
        'October': 'Octubre', 'November': 'Noviembre', 'December': 'Diciembre'
    }
    for eng, esp in meses.items():
        fecha_curso = fecha_curso.replace(eng, esp)

    # --- Contenido del certificado ---
    y_pos = height - 80

    y_pos -= 60
    draw_centered_text(c, "Se otorga el presente", y_pos, "Helvetica", 16, colors.HexColor('#4A90E2'))

    y_pos -= 50
    draw_centered_text(c, "RECONOCIMIENTO", y_pos, "Helvetica-Bold", 36)

    y_pos -= 60
    draw_centered_text(c, "A:", y_pos, "Helvetica-Bold", 24)

    y_pos -= 50
    draw_centered_text(c, nombre_completo.upper(), y_pos, "Helvetica-Bold", 20)

    y_pos -= 50
    c.setFont("Helvetica", 14)
    c.setFillColor(colors.black)
    texto_linea1 = f"Por su participación y gran desempeño en el curso de {nombre_curso.upper()},"
    texto_linea2 = "a su vez por aprobar el examen de conocimientos generales en el aula"
    texto_linea3 = f"con una duración de {duracion} horas en Pachuca Hidalgo."
    c.drawCentredString(width / 2, y_pos, texto_linea1)
    y_pos -= 25
    c.drawCentredString(width / 2, y_pos, texto_linea2)
    y_pos -= 25
    c.drawCentredString(width / 2, y_pos, texto_linea3)

    y_pos -= 40
    draw_centered_text(c, f"El día {fecha_curso} en Pachuca Hidalgo.", y_pos, "Helvetica-Bold", 14)

    # --- Firma ---
    y_pos -= 40
    firma_x = width / 2
    firma_y = y_pos
    draw_centered_text(c, "KARLA YISSEL PADILLA CADENA", firma_y, "Helvetica-Bold", 12)
    c.setStrokeColor(colors.black)
    c.setLineWidth(1)
    draw_centered_text(c, "CAPACITADOR", firma_y - 30, "Helvetica-Bold", 10)

    # Finalizar PDF
    c.showPage()
    c.save()
    buffer.seek(0)

    filename = f"reconocimiento_{request.user.username}_{curso.id}.pdf"
    progreso.certificado.save(filename, ContentFile(buffer.read()))
    progreso.fecha_aprobacion = timezone.now()
    progreso.save()

    buffer.seek(0)
    return FileResponse(progreso.certificado.open(), as_attachment=True, filename=filename)


@login_required
@user_passes_test(lambda u: u.is_superuser)
def crear_editar_examen(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)
    examen, _ = Examen.objects.get_or_create(curso=curso)

    if request.method == 'POST':
        examen_form = ExamenForm(request.POST, instance=examen)
        pregunta_formset = PreguntaFormSet(request.POST, instance=examen)

        if examen_form.is_valid() and pregunta_formset.is_valid():
            examen_form.save()
            preguntas = pregunta_formset.save()

            # Procesar opciones por cada pregunta
            for pregunta in preguntas:
                opciones_prefix = f'opciones_{pregunta.id}'
                opciones_formset = OpcionFormSet(
                    request.POST, instance=pregunta, prefix=opciones_prefix)
                if opciones_formset.is_valid():
                    opciones_formset.save()

            return redirect('detalle_curso', pk=curso_id)

    else:
        examen_form = ExamenForm(instance=examen)
        pregunta_formset = PreguntaFormSet(instance=examen)

    # Preparar formsets de opciones por pregunta
    opciones_formsets = []
    for pregunta_form in pregunta_formset.forms:
        pregunta = pregunta_form.instance
        prefix = f'opciones_{pregunta.id or "new"}'
        formset = OpcionFormSet(instance=pregunta, prefix=prefix)
        opciones_formsets.append((pregunta_form, formset))

    return render(request, 'crear_editar_examen.html', {
        'curso': curso,
        'examen_form': examen_form,
        'pregunta_formset': pregunta_formset,
        'opciones_formsets': opciones_formsets
    })


@login_required
@user_passes_test(lambda u: u.is_superuser)
def agregar_pregunta(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)
    examen, _ = Examen.objects.get_or_create(curso=curso)

    if request.method == 'POST':
        pregunta_form = PreguntaExamenForm(request.POST)
        formset = OpcionRespuestaFormSet(request.POST)

        if pregunta_form.is_valid() and formset.is_valid():
            pregunta = pregunta_form.save(commit=False)
            pregunta.examen = examen
            pregunta.save()
            formset.instance = pregunta
            formset.save()
            return redirect('detalle_curso', pk=curso_id)
    else:
        pregunta_form = PreguntaExamenForm()
        formset = OpcionRespuestaFormSet()

    return render(request, 'agregar_pregunta.html', {
        'pregunta_form': pregunta_form,
        'formset': formset,
        'curso_id': curso_id,
        'curso': curso,
    })


@login_required
@user_passes_test(lambda u: u.is_superuser)
def configurar_examen(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)
    examen, created = Examen.objects.get_or_create(curso=curso)

    if request.method == 'POST':
        tiempo = request.POST.get('tiempo_minutos')
        if tiempo and tiempo.isdigit():
            examen.tiempo_minutos = int(tiempo)
            examen.save()
            return redirect('detalle_curso', pk=curso_id)

    return render(request, 'configurar_examen.html', {
        'curso': curso,
        'examen': examen,
    })


@login_required
def listar_preguntas(request, curso_id):
    curso = get_object_or_404(Curso, id=curso_id)
    examen = getattr(curso, 'examen', None)
    if not examen:
        messages.error(request, "Este curso no tiene examen asignado.")
        return redirect('detalle_curso', pk=curso_id)

    preguntas = examen.preguntas.all()

    return render(request, 'listar_preguntas.html', {
        'curso': curso,
        'examen': examen,
        'preguntas': preguntas,
    })


def editar_pregunta(request, curso_id, pregunta_id):
    curso = get_object_or_404(Curso, id=curso_id)
    pregunta = get_object_or_404(
        PreguntaExamen, id=pregunta_id, examen__curso=curso)

    PreguntaForm = modelform_factory(PreguntaExamen, fields=['texto'])
    OpcionFormSet = inlineformset_factory(PreguntaExamen, OpcionRespuesta,
                                          fields=['texto', 'es_correcta'],
                                          extra=1, can_delete=True)

    if request.method == 'POST':
        pregunta_form = PreguntaForm(request.POST, instance=pregunta)
        opcion_formset = OpcionFormSet(request.POST, instance=pregunta)
        if pregunta_form.is_valid() and opcion_formset.is_valid():
            pregunta_form.save()
            opcion_formset.save()
            return redirect('listar_preguntas', curso_id=curso.id)
    else:
        pregunta_form = PreguntaForm(instance=pregunta)
        opcion_formset = OpcionFormSet(instance=pregunta)

    return render(request, 'editar_pregunta.html', {
        'curso': curso,
        'pregunta_form': pregunta_form,
        'formset': opcion_formset,
        'pregunta': pregunta,
    })


@login_required
@user_passes_test(lambda u: u.is_superuser)
def eliminar_pregunta(request, pregunta_id):
    pregunta = get_object_or_404(PreguntaExamen, id=pregunta_id)
    curso_id = pregunta.examen.curso.id
    pregunta.delete()
    messages.success(request, "Pregunta eliminada correctamente.")
    return redirect('listar_preguntas', curso_id=curso_id)


def es_superusuario(user):
    return user.is_superuser


@user_passes_test(lambda u: u.is_superuser)
def usuarios_aprobados_lista(request):
    progresos = ProgresoCurso.objects.filter(
        examen_aprobado=True).select_related('usuario', 'curso')
    datos = []

    for progreso in progresos:
        perfil = PerfilUsuario.objects.filter(user=progreso.usuario).first()
        datos.append({
            'usuario': progreso.usuario,
            'perfil': perfil,
            'curso': progreso.curso,
            'progreso': progreso,
        })

    return render(request, 'usuarios_aprobados_lista.html', {
        'usuarios_aprobados': datos
    })
